#!/usr/bin/env python3
"""
PPT 多智能体自动生成流水线
Pipeline: Agent-A(风格判定) → Agent-B(视觉采集) → Agent-C(内容融合) → 监督评审
"""

import json
import os
from pathlib import Path
from dataclasses import dataclass, field, asdict
from typing import Optional


# ─────────────────────────────────────────────
# 数据结构定义
# ─────────────────────────────────────────────

@dataclass
class UserRequirement:
    """用户需求（管家收集）"""
    topic: str
    audience: str
    slide_count: int = 12
    content_outline: str = ""
    output_format: str = "pptx"  # pptx | html | markdown
    special_notes: str = ""
    brand_colors: Optional[dict] = None


@dataclass
class StyleProfile:
    """Agent-A 输出：设计规范"""
    style_name: str
    rationale: str
    colors: dict = field(default_factory=dict)
    typography: dict = field(default_factory=dict)
    layout_templates: list = field(default_factory=list)
    slide_aspect: str = "16:9"
    icon_style: str = "line"


@dataclass
class VisualAssets:
    """Agent-B 输出：视觉素材清单"""
    slides_assets: list = field(default_factory=list)
    color_palette_applied: bool = False
    total_assets: int = 0


@dataclass
class PresentationOutput:
    """Agent-C 输出：最终演示文稿"""
    title: str
    total_slides: int
    output_path: str
    slides: list = field(default_factory=list)
    style_applied: str = ""


# ─────────────────────────────────────────────
# 管家：需求收集 & 任务调度
# ─────────────────────────────────────────────

class ButlerAgent:
    """管家智能体：用户唯一接口，负责需求收集与流水线调度"""

    AUDIENCE_MAP = {
        "高管": "极简商务",
        "董事会": "极简商务",
        "客户": "品牌商务",
        "投资人": "科技炫彩",
        "技术团队": "工程简洁",
        "学生": "活泼教育",
        "公众": "活泼教育",
        "学术": "学术专业",
        "教授": "学术专业",
    }

    def collect_requirements(self, user_input: dict) -> UserRequirement:
        """从用户输入中提取需求，缺失字段设为默认值"""
        req = UserRequirement(
            topic=user_input.get("topic", "未命名主题"),
            audience=user_input.get("audience", "通用受众"),
            slide_count=user_input.get("slide_count", 12),
            content_outline=user_input.get("content_outline", ""),
            output_format=user_input.get("output_format", "pptx"),
            special_notes=user_input.get("special_notes", ""),
            brand_colors=user_input.get("brand_colors", None),
        )
        return req

    def infer_style_hint(self, req: UserRequirement) -> str:
        """根据受众推断风格提示"""
        for keyword, style in self.AUDIENCE_MAP.items():
            if keyword in req.audience:
                return style
        return "极简商务"

    def run_pipeline(self, user_input: dict) -> PresentationOutput:
        """执行完整的多智能体流水线"""
        print("【管家】开始收集需求...")
        req = self.collect_requirements(user_input)
        style_hint = self.infer_style_hint(req)

        print(f"【管家】已分析受众「{req.audience}」，推荐风格：{style_hint}")
        print(f"【管家】→ 触发 Agent-A（风格判定）")

        # 阶段1：风格判定
        agent_a = StyleDeterminationAgent()
        style_profile = agent_a.execute(req, style_hint)
        print(f"【管家】Agent-A 完成，风格：{style_profile.style_name}")

        print(f"【管家】→ 触发 Agent-B（视觉采集）")

        # 阶段2：视觉采集
        agent_b = VisualCollectionAgent()
        visual_assets = agent_b.execute(req, style_profile)
        print(f"【管家】Agent-B 完成，采集素材 {visual_assets.total_assets} 个")

        print(f"【管家】→ 触发 Agent-C（内容融合）")

        # 阶段3：内容融合
        agent_c = ContentFusionAgent()
        result = agent_c.execute(req, style_profile, visual_assets)
        print(f"【管家】Agent-C 完成，生成 {result.total_slides} 页幻灯片")

        # 监督评审
        supervisor = SupervisorAgent()
        passed, feedback = supervisor.review(result, style_profile, visual_assets)

        retry_count = 0
        max_retry = 3
        while not passed and retry_count < max_retry:
            retry_count += 1
            print(f"【监督】评审未通过，退回重做（第{retry_count}轮）")
            print(f"【监督】问题清单：{feedback}")

            # 根据反馈重新执行相关 Agent
            if feedback.get("rework_agent_b"):
                visual_assets = agent_b.execute(req, style_profile)

            result = agent_c.execute(req, style_profile, visual_assets,
                                     feedback=feedback)
            passed, feedback = supervisor.review(result, style_profile, visual_assets)

        if passed:
            print(f"【监督】✅ 评审通过！")
            print(f"【管家】任务完成，PPT 已生成：{result.output_path}")
        else:
            print(f"【管家】❌ 已达最大重试次数({max_retry})，请用户决策：")
            print("  1. 降标使用当前版本")
            print("  2. 补充信息重试")
            print("  3. 放弃")

        return result


# ─────────────────────────────────────────────
# Agent-A：风格判定智能体
# ─────────────────────────────────────────────

class StyleDeterminationAgent:
    """风格判定智能体：分析受众，输出设计规范"""

    STYLE_PRESETS = {
        "极简商务": {
            "colors": {
                "primary": "#1A1A2E",
                "secondary": "#16213E",
                "accent": "#0F3460",
                "highlight": "#E94560",
                "background": "#FFFFFF",
                "text_dark": "#1A1A1A",
                "text_light": "#666666",
            },
            "typography": {
                "title": {"font": "思源黑体 Bold", "size_pt": 40, "weight": "bold"},
                "subtitle": {"font": "思源黑体 Medium", "size_pt": 28, "weight": "medium"},
                "body": {"font": "思源宋体 Regular", "size_pt": 18, "weight": "regular"},
                "caption": {"font": "思源黑体 Light", "size_pt": 14, "weight": "light"},
            },
            "layout_templates": [
                "cover_centered", "toc_numbered", "content_two_column",
                "data_chart_full", "quote_emphasis", "closing_cta"
            ],
            "icon_style": "line",
        },
        "学术专业": {
            "colors": {
                "primary": "#2C3E50", "secondary": "#34495E",
                "accent": "#2980B9", "highlight": "#27AE60",
                "background": "#FAFAFA", "text_dark": "#2C3E50",
                "text_light": "#7F8C8D",
            },
            "typography": {
                "title": {"font": "思源黑体 Medium", "size_pt": 36, "weight": "medium"},
                "subtitle": {"font": "思源黑体 Regular", "size_pt": 26, "weight": "regular"},
                "body": {"font": "思源宋体 Regular", "size_pt": 18, "weight": "regular"},
                "caption": {"font": "思源黑体 Light", "size_pt": 13, "weight": "light"},
            },
            "layout_templates": [
                "cover_academic", "abstract_page", "content_standard",
                "data_research", "conclusion_page", "reference_list"
            ],
            "icon_style": "filled",
        },
        "科技炫彩": {
            "colors": {
                "primary": "#0D0D0D", "secondary": "#1A1A1A",
                "accent": "#00D4FF", "highlight": "#FF6B35",
                "background": "#050505", "text_dark": "#FFFFFF",
                "text_light": "#AAAAAA",
            },
            "typography": {
                "title": {"font": "Inter Bold", "size_pt": 44, "weight": "bold"},
                "subtitle": {"font": "Inter Medium", "size_pt": 30, "weight": "medium"},
                "body": {"font": "Inter Regular", "size_pt": 18, "weight": "regular"},
                "caption": {"font": "Inter Light", "size_pt": 13, "weight": "light"},
            },
            "layout_templates": [
                "cover_hero", "feature_spotlight", "metric_dashboard",
                "comparison_side", "roadmap_timeline", "closing_impact"
            ],
            "icon_style": "gradient",
        },
        "活泼教育": {
            "colors": {
                "primary": "#FF6B6B", "secondary": "#4ECDC4",
                "accent": "#45B7D1", "highlight": "#FFA07A",
                "background": "#FFFEF7", "text_dark": "#333333",
                "text_light": "#888888",
            },
            "typography": {
                "title": {"font": "圆体 Bold", "size_pt": 42, "weight": "bold"},
                "subtitle": {"font": "圆体 Medium", "size_pt": 28, "weight": "medium"},
                "body": {"font": "黑体 Regular", "size_pt": 18, "weight": "regular"},
                "caption": {"font": "黑体 Light", "size_pt": 14, "weight": "light"},
            },
            "layout_templates": [
                "cover_colorful", "agenda_visual", "content_card",
                "exercise_page", "summary_checklist", "closing_fun"
            ],
            "icon_style": "colored",
        },
        "创意自由": {
            "colors": {
                "primary": "#6C5CE7", "secondary": "#A29BFE",
                "accent": "#FD79A8", "highlight": "#FDCB6E",
                "background": "#FFFFFF", "text_dark": "#2D3436",
                "text_light": "#636E72",
            },
            "typography": {
                "title": {"font": "展示字体 Bold", "size_pt": 48, "weight": "bold"},
                "subtitle": {"font": "无衬线 Medium", "size_pt": 30, "weight": "medium"},
                "body": {"font": "无衬线 Regular", "size_pt": 18, "weight": "regular"},
                "caption": {"font": "无衬线 Light", "size_pt": 13, "weight": "light"},
            },
            "layout_templates": [
                "cover_asymmetric", "moodboard_page", "content_creative",
                "showcase_gallery", "process_visual", "closing_creative"
            ],
            "icon_style": "illustrated",
        },
    }

    def execute(self, req: UserRequirement, style_hint: str) -> StyleProfile:
        """分析需求，输出设计规范"""
        # 使用用户品牌色覆盖预设色
        preset = self.STYLE_PRESETS.get(style_hint, self.STYLE_PRESETS["极简商务"])
        colors = preset["colors"].copy()

        if req.brand_colors:
            colors.update(req.brand_colors)

        return StyleProfile(
            style_name=style_hint,
            rationale=f"受众「{req.audience}」，主题「{req.topic}」，匹配{style_hint}风格",
            colors=colors,
            typography=preset["typography"],
            layout_templates=preset["layout_templates"],
            slide_aspect="16:9",
            icon_style=preset["icon_style"],
        )


# ─────────────────────────────────────────────
# Agent-B：视觉采集智能体
# ─────────────────────────────────────────────

class VisualCollectionAgent:
    """视觉采集智能体：根据内容语义采集和优化视觉素材"""

    SLIDE_TYPE_MAP = {
        "cover": "background",
        "toc": "icon",
        "content": "illustration",
        "data": "chart",
        "comparison": "table",
        "quote": "portrait",
        "closing": "brand",
    }

    def execute(self, req: UserRequirement,
                style_profile: StyleProfile) -> VisualAssets:
        """为每张幻灯片分析并规划视觉素材"""
        slides_assets = []
        slide_count = req.slide_count

        # 生成每页素材规划
        for i in range(slide_count):
            slide_type = self._infer_slide_type(i, slide_count)
            asset_plan = self._plan_asset(i + 1, slide_type,
                                          req.topic, style_profile)
            slides_assets.append(asset_plan)

        total = sum(len(s.get("assets", [])) for s in slides_assets)

        return VisualAssets(
            slides_assets=slides_assets,
            color_palette_applied=True,
            total_assets=total,
        )

    def _infer_slide_type(self, index: int, total: int) -> str:
        if index == 0:
            return "cover"
        elif index == 1:
            return "toc"
        elif index == total - 1:
            return "closing"
        elif index % 4 == 0:
            return "data"
        else:
            return "content"

    def _plan_asset(self, slide_idx: int, slide_type: str,
                    topic: str, style: StyleProfile) -> dict:
        visual_role = self.SLIDE_TYPE_MAP.get(slide_type, "illustration")
        return {
            "slide_index": slide_idx,
            "slide_type": slide_type,
            "assets": [
                {
                    "role": visual_role,
                    "source": "generated",
                    "search_query": f"{topic} {slide_type} {style.style_name}",
                    "color_filter": style.colors.get("accent", "#000000"),
                    "overlay_opacity": 0.35 if slide_type == "cover" else 0,
                    "placeholder": f"[{slide_type}_{visual_role}_{slide_idx}]",
                }
            ],
        }


# ─────────────────────────────────────────────
# Agent-C：内容融合智能体
# ─────────────────────────────────────────────

class ContentFusionAgent:
    """内容融合智能体：整合文本、图像、数据，生成最终幻灯片结构"""

    MAX_BULLET_POINTS = 5
    MAX_WORDS_PER_BULLET = 20

    def execute(self, req: UserRequirement,
                style_profile: StyleProfile,
                visual_assets: VisualAssets,
                feedback: Optional[dict] = None) -> PresentationOutput:
        """融合所有输入，生成演示文稿结构"""
        slides = []

        for i in range(req.slide_count):
            slide_asset = (visual_assets.slides_assets[i]
                           if i < len(visual_assets.slides_assets) else {})
            slide = self._build_slide(i, req, style_profile, slide_asset, feedback)
            slides.append(slide)

        output_dir = Path("output")
        output_dir.mkdir(exist_ok=True)
        output_path = str(output_dir / f"{req.topic.replace(' ', '_')}.pptx")

        # 调用 python-pptx 生成文件（骨架，AI 执行时会替换为实际调用）
        self._generate_pptx(slides, style_profile, output_path)

        return PresentationOutput(
            title=req.topic,
            total_slides=len(slides),
            output_path=output_path,
            slides=slides,
            style_applied=style_profile.style_name,
        )

    def _build_slide(self, idx: int, req: UserRequirement,
                     style: StyleProfile, asset: dict,
                     feedback: Optional[dict]) -> dict:
        slide_type = asset.get("slide_type", "content")
        layout = self._select_layout(slide_type, style)
        return {
            "index": idx + 1,
            "type": slide_type,
            "layout": layout,
            "title": f"[第{idx+1}页标题]",
            "body": [],
            "assets": asset.get("assets", []),
            "colors": style.colors,
            "typography": style.typography,
        }

    def _select_layout(self, slide_type: str, style: StyleProfile) -> str:
        templates = style.layout_templates
        mapping = {
            "cover": 0, "toc": 1, "content": 2,
            "data": 3, "quote": 4, "closing": 5,
        }
        idx = mapping.get(slide_type, 2)
        return templates[idx] if idx < len(templates) else templates[0]

    def _generate_pptx(self, slides: list, style: StyleProfile,
                       output_path: str) -> None:
        """
        实际 PPT 生成逻辑骨架。
        AI 执行时使用 python-pptx 替换占位符生成真实文件。
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            slide_layout = prs.slide_layouts[6]  # Blank layout

            for slide_data in slides:
                slide = prs.slides.add_slide(slide_layout)
                title_shape = slide.shapes.add_textbox(
                    Inches(1), Inches(1), Inches(11.33), Inches(1.5)
                )
                tf = title_shape.text_frame
                tf.text = slide_data.get("title", "")

                # 应用字体样式
                title_font = style.typography.get("title", {})
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(title_font.get("size_pt", 40))
                        hex_color = style.colors.get("text_dark", "#1A1A1A").lstrip("#")
                        run.font.color.rgb = RGBColor(
                            int(hex_color[0:2], 16),
                            int(hex_color[2:4], 16),
                            int(hex_color[4:6], 16)
                        )

            prs.save(output_path)
            print(f"  ✓ PPT 已保存至：{output_path}")

        except ImportError:
            # python-pptx 未安装时，输出 JSON 结构文件作为替代
            json_path = output_path.replace(".pptx", "_structure.json")
            with open(json_path, "w", encoding="utf-8") as f:
                json.dump({
                    "slides": slides,
                    "style": asdict(style) if hasattr(style, '__dataclass_fields__') else str(style),
                    "note": "安装 python-pptx 后可生成真实 .pptx 文件"
                }, f, ensure_ascii=False, indent=2)
            print(f"  ⚠ python-pptx 未安装，已输出结构文件：{json_path}")


# ─────────────────────────────────────────────
# 监督：质量评审智能体
# ─────────────────────────────────────────────

class SupervisorAgent:
    """监督智能体：按五维度评审最终输出"""

    PASS_THRESHOLD = 3

    def review(self, result: PresentationOutput,
               style_profile: StyleProfile,
               visual_assets: VisualAssets) -> tuple[bool, dict]:
        """评审演示文稿质量"""
        scores = {}

        # 完整性：检查幻灯片数量
        scores["completeness"] = 4 if result.total_slides >= 10 else 2

        # 正确性：检查结构完整性
        scores["correctness"] = 4 if all(
            s.get("layout") and s.get("title") for s in result.slides
        ) else 2

        # 规范性：检查是否应用了风格
        scores["conformance"] = 4 if result.style_applied == style_profile.style_name else 2

        # 可用性：检查输出文件存在
        scores["usability"] = 4 if Path(result.output_path).exists() or \
            Path(result.output_path.replace(".pptx", "_structure.json")).exists() else 2

        # 美观度：检查素材覆盖率
        coverage = visual_assets.total_assets / max(result.total_slides, 1)
        scores["aesthetics"] = 4 if coverage >= 0.8 else (3 if coverage >= 0.5 else 2)

        # 判断是否通过
        passed = all(v >= self.PASS_THRESHOLD for v in scores.values())

        feedback = {
            "scores": scores,
            "passed": passed,
            "rework_agent_b": scores.get("aesthetics", 0) < self.PASS_THRESHOLD,
            "rework_agent_c": scores.get("conformance", 0) < self.PASS_THRESHOLD or
                               scores.get("usability", 0) < self.PASS_THRESHOLD,
        }

        self._print_review(scores, passed, feedback)
        return passed, feedback

    def _print_review(self, scores: dict, passed: bool, feedback: dict):
        status = "✅ 通过" if passed else "❌ 退回"
        print(f"\n【监督评审结论】：{status}")
        print(f"  完整性： {'★' * scores.get('completeness', 0)} {scores.get('completeness', 0)}分")
        print(f"  正确性： {'★' * scores.get('correctness', 0)} {scores.get('correctness', 0)}分")
        print(f"  规范性： {'★' * scores.get('conformance', 0)} {scores.get('conformance', 0)}分")
        print(f"  可用性： {'★' * scores.get('usability', 0)} {scores.get('usability', 0)}分")
        print(f"  美观度： {'★' * scores.get('aesthetics', 0)} {scores.get('aesthetics', 0)}分")


# ─────────────────────────────────────────────
# 入口
# ─────────────────────────────────────────────

if __name__ == "__main__":
    # 示例：生成一份销售复盘 PPT
    user_input = {
        "topic": "2025年Q3销售复盘",
        "audience": "高管 销售总监",
        "slide_count": 12,
        "content_outline": "市场概况、销售数据、Top案例、问题分析、Q4规划",
        "output_format": "pptx",
        "special_notes": "需要包含数据图表",
    }

    butler = ButlerAgent()
    result = butler.run_pipeline(user_input)

    print(f"\n✅ 演示文稿生成完成！")
    print(f"   标题：{result.title}")
    print(f"   页数：{result.total_slides}")
    print(f"   路径：{result.output_path}")
    print(f"   风格：{result.style_applied}")
