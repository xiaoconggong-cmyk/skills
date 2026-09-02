"""Generate demo PPTX using v5 strict templates."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ===== Constants =====
SLIDE_W, SLIDE_H = 10.00, 7.50
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_RED   = RGBColor(0xFF, 0x00, 0x00)
C_PURPLE= RGBColor(0x70, 0x30, 0xA0)
C_DASH  = RGBColor(0x9C, 0xA3, 0xAF)
S_COVER=Pt(40); S_TOC=Pt(32); S_ITEM=Pt(24); S_H=Pt(16); S_TAG=Pt(14); S_TBL=Pt(12)
AX, AY = 1.39, 0.19  # section title anchor
SX, SY = 0.00, 0.78  # subtitle anchor
BX, BY = 0.14, 1.29  # body anchor
TY, TLX, TRX, TW, TH = 4.41, 5.62, 7.59, 1.33, 0.93  # tag anchors

def dash_box(slide, x, y, w, h):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.background()
    s.line.color.rgb = C_DASH
    s.line.width = Pt(0.5)
    s.line.dash_style = 2

def tb_bold(slide, text, x, y, w, h, size, color, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = True; r.font.color.rgb = color

def body_tb(slide, paras, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for pi, pr in enumerate(paras):
        p = tf.paragraphs[pi] if pi == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        for rd in pr:
            if rd[0] is None: continue
            r = p.add_run()
            r.text = rd[0]
            r.font.size = S_H
            r.font.bold = rd[1] if len(rd) > 1 else False
            r.font.color.rgb = rd[2] if len(rd) > 2 else C_BLACK

def add_tag(slide, text, x, y):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(TW), Inches(TH))
    s.fill.background()
    s.line.fill.background()
    p = s.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = S_TAG; r.font.bold = True; r.font.color.rgb = C_PURPLE

# ===== Generate =====
prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)

# --- Slide 1: Cover (Template A) ---
s = prs.slides.add_slide(prs.slide_layouts[6])
tb_bold(s, "深度学习异常检测研究进展", 0.42, 2.41, 9.16, 1.61, S_COVER, C_BLACK, PP_ALIGN.CENTER)
tb_bold(s, "主动视觉语言推理方法", 0.50, 4.47, 9.00, 0.40, S_H, C_BLACK, PP_ALIGN.CENTER)
tb_bold(s, "2026年3月  大组会汇报", 0.50, 4.90, 9.00, 0.40, S_H, C_BLACK, PP_ALIGN.CENTER)

# --- Slide 2: TOC (Template B) ---
s = prs.slides.add_slide(prs.slide_layouts[6])
tb_bold(s, "目录", 4.00, 0.07, 1.09, 0.64, S_TOC, C_BLACK)
items = ["1. 项目背景", "2. 研究现状", "3. 问题与挑战", "4. 工作进展", "5. 项目落地", "6. 后续计划"]
tb = s.shapes.add_textbox(Inches(2.48), Inches(1.43), Inches(4.77), Inches(3.98))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "\n".join(items); r.font.size = S_ITEM; r.font.color.rgb = C_BLACK

# --- Slide 3: Content page (Template C) ---
s = prs.slides.add_slide(prs.slide_layouts[6])
tb_bold(s, "1. 项目背景", AX, AY, 5.00, 0.40, S_H, C_BLACK)
tb_bold(s, "主动视觉语言推理", SX, SY, 10.00, 0.40, S_H, C_BLACK)
dash_box(s, 0.10, 1.25, 9.60, 1.70)
body_tb(s, [
    [("在真实机器人应用场景中（如工业巡检、物体抓取），", False, C_BLACK),
     ("关键信息", True, C_RED),
     ("往往仅分布于局部区域，单次被动感知难以支撑稳定可靠判断。", False, C_BLACK)],
    [(None,)],
    [("机器人不应被动地均匀理解场景中每一个像素，而应", False, C_BLACK),
     ("主动选择", True, C_RED),
     ("对当前任务最有价值的信息区域进行聚焦感知与推理。", False, C_BLACK)],
], BX, BY, 9.55, 1.70)
add_tag(s, "被动感知", TLX, TY)
add_tag(s, "主动聚焦", TRX, TY)

# --- Slide 4: Content page 2 (Template C) ---
s = prs.slides.add_slide(prs.slide_layouts[6])
tb_bold(s, "3. 问题与挑战", AX, AY, 5.00, 0.40, S_H, C_BLACK)
tb_bold(s, "在复杂的工业环境中，如何使机器人能够主动选择关键信息？", SX, SY, 10.00, 0.90, S_H, C_BLACK)
dash_box(s, 0.10, 1.60, 9.60, 1.50)
body_tb(s, [
    [("主动关键场景选择能力：", True, C_BLACK),
     ("在复杂环境中优先关注对任务决策最具价值的场景数据。", False, C_BLACK)],
    [("主动语义歧义建模能力：", True, C_BLACK),
     ("针对语言与视觉之间的模糊对应关系，强化多模态对齐与理解。", False, C_BLACK)],
    [("主动聚焦与推理能力：", True, C_BLACK),
     ("聚焦语义最关键的区域，逐步递归推理得到精确判断。", False, C_BLACK)],
], BX, 1.65, 9.55, 1.50)
add_tag(s, "语义歧义", TLX, 5.00)
add_tag(s, "递归推理", TRX, 5.00)

# --- Slide 5: Table page (Template E) ---
s = prs.slides.add_slide(prs.slide_layouts[6])
tb_bold(s, "4. 实验结果", AX, AY, 5.00, 0.40, S_H, C_BLACK)
headers = ["Method", "MVTec", "VisA", "BTAD", "Average"]
rows = [
    ["PaDiM", "85.7", "72.1", "78.4", "78.7"],
    ["PatchCore", "87.3", "74.6", "80.2", "80.7"],
    ["CFA", "88.9", "76.3", "82.1", "82.4"],
    ["Ours", "94.2", "88.5", "90.3", "91.0"],
]
n_rows = len(rows) + 1; n_cols = len(headers)
ts = s.shapes.add_table(n_rows, n_cols, Inches(0.15), Inches(1.15), Inches(5.00), Inches(n_rows * 0.40))
t = ts.table
for ci, h in enumerate(headers):
    c = t.cell(0, ci); c.text = ""
    p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = h; r.font.size = S_TBL; r.font.bold = True; r.font.color.rgb = C_BLACK
for ri, row in enumerate(rows):
    for ci, v in enumerate(row):
        c = t.cell(ri + 1, ci); c.text = ""
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = v; r.font.size = S_TBL
        is_ours = (ri == len(rows) - 1)
        r.font.bold = is_ours
        r.font.color.rgb = C_RED if is_ours else C_BLACK
dash_box(s, 0.10, 1.10, 5.10, n_rows * 0.40 + 0.20)
dash_box(s, 5.31, 1.10, 4.24, 1.50)
tb_bold(s, "关键发现", 5.36, 1.15, 4.14, 0.30, S_H, C_BLACK)
tb2 = s.shapes.add_textbox(Inches(5.36), Inches(1.50), Inches(4.14), Inches(1.50))
tf2 = tb2.text_frame; tf2.word_wrap = True
for i, line in enumerate(["  MVTec-AD 超越 SOTA 5.3%", "  VisA 提升 12.2%", "  推理速度 23 FPS"]):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    r = p.add_run(); r.text = line; r.font.size = Pt(14); r.font.color.rgb = C_BLACK

# --- Slide 6: Summary page (Template F) ---
s = prs.slides.add_slide(prs.slide_layouts[6])
tb_bold(s, "6. 后续计划", AX, AY, 5.00, 0.40, S_H, C_BLACK)
sum_items = [
    "将生成的缺陷图片（200+张）加入训练数据，重新微调缺陷检测模型",
    "AI验布机重点解决类别不均衡问题",
    "撰写毕业论文和相关专利"
]
dash_box(s, 0.72, 1.03, 8.62, len(sum_items) * 0.30 + 0.20)
for i, item in enumerate(sum_items):
    tb_bold(s, f"{i+1}. {item}", 0.77, 1.08 + i * 0.30, 8.52, 0.30, S_H, C_BLACK)

out = r"C:\Users\25124\.workbuddy\skills\ppt-multiagent\assets\demo_v5.pptx"
prs.save(out)
print(f"Demo saved: {out}, {len(prs.slides)} slides")
