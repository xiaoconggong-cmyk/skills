"""Screenshot v5 demo PPTX using Playwright."""
import io, os, sys
sys.path.insert(0, r"C:\Users\25124\.workbuddy\skills\ppt-multiagent\assets")

from pptx import Presentation
from playwright.sync_api import sync_playwright
import http.server, socketserver, threading

ASSETS_DIR = r"C:\Users\25124\.workbuddy\skills\ppt-multiagent\assets"
SCREENSHOTS_DIR = os.path.join(ASSETS_DIR, "screenshots_v5")
os.makedirs(SCREENSHOTS_DIR, exist_ok=True)

# 1. Convert PPTX slides to HTML preview
pptx_path = os.path.join(ASSETS_DIR, "demo_v5.pptx")
prs = Presentation(pptx_path)

html_parts = []
html_parts.append('''<!DOCTYPE html><html lang="zh"><head><meta charset="UTF-8">
<style>
* { margin: 0; padding: 0; box-sizing: border-box; }
body { background: #e5e7eb; padding: 20px; font-family: sans-serif; }
h1 { text-align: center; color: #333; margin-bottom: 20px; }
.slide-container { margin-bottom: 30px; background: white; box-shadow: 0 2px 8px rgba(0,0,0,0.15); }
.slide-header { background: #f3f4f6; padding: 8px 16px; font-size: 14px; color: #666; border-bottom: 1px solid #e5e7eb; }
.slide-canvas { position: relative; background: white; overflow: hidden; }
</style></head><body><h1>PPT v5 Demo - 模板效果预览</h1>''')

for idx, slide in enumerate(prs.slides):
    html_parts.append(f'<div class="slide-container">')
    html_parts.append(f'<div class="slide-header">第 {idx+1} 页</div>')
    
    # Render at 96 DPI
    w_px = int(10.00 * 96)
    h_px = int(7.50 * 96)
    
    html_parts.append(f'<div class="slide-canvas" style="width:{w_px}px;height:{h_px}px;">')
    
    for shape in slide.shapes:
        left = int(shape.left / 914400 * 96)
        top = int(shape.top / 914400 * 96)
        w = int(shape.width / 914400 * 96)
        h = int(shape.height / 914400 * 96)
        
        text_content = ""
        if shape.has_text_frame:
            text_content = shape.text_frame.text
            # Get font details
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    f = run.font
                    size_px = int(f.size / 12700 * 1.333) if f.size else 16
                    is_bold = 'font-weight:bold;' if f.bold else ''
                    color = '#000000'
                    try:
                        if f.color and f.color.type is not None:
                            try: color = f'#{f.color.rgb}'
                            except: pass
                    except: pass
                    
                    html_parts.append(
                        f'<div style="position:absolute;left:{left}px;top:{top}px;'
                        f'width:{w}px;height:{h}px;font-size:{size_px}px;{is_bold}'
                        f'color:{color};overflow:hidden;word-wrap:break-word;">'
                        f'{run.text}</div>'
                    )
        
        # Check for table
        if shape.has_table:
            html_parts.append(f'<div style="position:absolute;left:{left}px;top:{top}px;'
                            f'width:{w}px;height:{h}px;border:1px solid #ccc;font-size:12px;">')
            tbl = shape.table
            for ri in range(tbl.rows.__len__()):
                row_h = h // max(tbl.rows.__len__(), 1)
                for ci in range(len(tbl.columns)):
                    col_w = w // max(len(tbl.columns), 1)
                    cell = tbl.cell(ri, ci)
                    html_parts.append(
                        f'<div style="position:absolute;left:{ci*col_w}px;top:{ri*row_h}px;'
                        f'width:{col_w}px;height:{row_h}px;border:1px solid #ddd;'
                        f'display:flex;align-items:center;justify-content:center;'
                        f'text-align:center;">{cell.text}</div>'
                    )
            html_parts.append('</div>')
        
        # Rectangles (dashed boxes, shapes)
        if str(shape.shape_type) in ('AUTO_SHAPE (1)', 'RECTANGLE'):
            fill_color = 'none'
            stroke = '#9CA3AF'
            try:
                if shape.fill.type is not None:
                    try:
                        fill_color = f'#{shape.fill.fore_color.rgb}'
                    except:
                        fill_color = '#f0f0f0'
            except: pass
            
            html_parts.append(
                f'<div style="position:absolute;left:{left}px;top:{top}px;'
                f'width:{w}px;height:{h}px;border:1px dashed {stroke};'
                f'background:{fill_color};"></div>'
            )
    
    html_parts.append('</div></div>')

html_parts.append('</body></html>')

html_path = os.path.join(ASSETS_DIR, "demo_v5_preview.html")
with open(html_path, 'w', encoding='utf-8') as f:
    f.write('\n'.join(html_parts))
print(f"HTML preview: {html_path}")

# 2. Screenshot with Playwright
PORT = 8999

class Handler(http.server.SimpleHTTPRequestHandler):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, directory=ASSETS_DIR, **kwargs)

def run_server():
    with socketserver.TCPServer(("", PORT), Handler) as httpd:
        httpd.serve_forever()

server_thread = threading.Thread(target=run_server, daemon=True)
server_thread.start()

with sync_playwright() as p:
    browser = p.chromium.launch(headless=True)
    page = browser.new_page(viewport={"width": 1200, "height": 900})
    
    # Full preview screenshot
    page.goto(f"http://localhost:{PORT}/demo_v5_preview.html", timeout=15000)
    page.wait_for_timeout(1000)
    full_h = page.evaluate("document.body.scrollHeight")
    page.set_viewport_size({"width": 1200, "height": full_h})
    page.screenshot(path=os.path.join(SCREENSHOTS_DIR, "full_preview.png"), full_page=True)
    print(f"Screenshot: full_preview.png")
    
    browser.close()

print("All screenshots done!")
