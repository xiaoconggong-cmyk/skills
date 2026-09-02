# -*- coding: utf-8 -*-
"""v7.0 验证 Demo：测试溢出防护 + 空页增强 + 自检循环"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W, SLIDE_H = 10.00, 7.50
SAFE_LEFT, SAFE_RIGHT = 0.14, 9.69
SAFE_TOP, SAFE_BOTTOM = 0.19, 7.30
SAFE_WIDTH = SAFE_RIGHT - SAFE_LEFT
TITLE_X_POOL = [1.39, 1.35, 1.28, 1.20, 1.14]
PAGE_NUM_X, PAGE_NUM_Y = 9.40, 7.10

C_BLACK = RGBColor(0x00,0x00,0x00)
C_RED   = RGBColor(0xFF,0x00,0x00)
C_PURPLE= RGBColor(0x70,0x30,0xA0)
C_DASH  = RGBColor(0x9C,0xA3,0xAF)
S_HEADING, S_TAG, S_TABLE, S_PAGE_NUM = Pt(16), Pt(14), Pt(12), Pt(10)

def _txt(slide, text, x, y, w, h, size=S_HEADING, color=C_BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold; r.font.color.rgb = color
    return tb

def estimate_text_height(text, width_in, font_pt=16, line_h=1.15, space_after_pt=4):
    chars_per_line = int(width_in * 72 / (font_pt * 0.55))
    n_lines = max(1, -(-len(text) // chars_per_line))
    line_in = font_pt * line_h / 72.0
    return n_lines * line_in + space_after_pt / 72.0 + 0.10

def add_body_separate(slide, paras, x, start_y, width):
    """每要点独立 textbox，返回 end_y"""
    cy = start_y
    for para in paras:
        ph = estimate_text_height(para, width, 16)
        if cy + ph > SAFE_BOTTOM:
            break
        tb = slide.shapes.add_textbox(Inches(x), Inches(cy), Inches(width), Inches(ph+0.05))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.space_after = Pt(4)
        # 找加粗前缀
        if '：' in para[:20]:
            head, rest = para.split('：', 1)
            r1 = p.add_run(); r1.text = head + '：'; r1.font.size = S_HEADING; r1.font.bold = True; r1.font.color.rgb = C_BLACK
            r2 = p.add_run(); r2.text = rest; r2.font.size = S_HEADING; r2.font.color.rgb = C_BLACK
        else:
            r = p.add_run(); r.text = para; r.font.size = S_HEADING; r.font.color.rgb = C_BLACK
        cy += ph + 0.08
    return cy

def add_page_number(slide, n):
    _txt(slide, str(n), PAGE_NUM_X, PAGE_NUM_Y, 0.50, 0.25, S_PAGE_NUM, C_DASH, align=PP_ALIGN.RIGHT)

def add_dashed_box(slide, x, y, w, h):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.background(); sp.line.color.rgb = C_DASH; sp.line.width = Pt(0.5)
    sp.line.dash_style = 2
    return sp

def add_decoration_line(slide, x1, y1, x2, y2, color=C_DASH, w=0.5):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color; ln.line.width = Pt(w)
    return ln

def add_status_card(slide, x, y, w, h, title, body):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.background(); sp.line.color.rgb = C_DASH; sp.line.width = Pt(0.75)
    tf = sp.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title; r.font.size = S_TAG; r.font.bold = True; r.font.color.rgb = C_BLACK
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = body; r2.font.size = Pt(11); r2.font.color.rgb = C_BLACK
    return sp

def check_overflow(slide):
    """溢出检测：排除页脚页码（位于 7.30" 以下页脚区，合法存在）。"""
    viol = []
    for shape in slide.shapes:
        x = shape.left / 914400.0
        y = shape.top / 914400.0
        bottom = (shape.top + shape.height) / 914400.0
        # 页脚页码：右下角 (x>9.0 且 y>7.0) 豁免
        if x > 9.0 and y > 7.0:
            continue
        if bottom > SAFE_BOTTOM + 0.01:
            viol.append((shape.name, round(bottom, 2)))
    return viol

def build_cover(prs, n):
    """封面：焦点留白 + 装饰线 + 锚点（解决 v6.1 上半页空洞）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # 大标题居中偏上
    _txt(s, "第二周工作周报", 2.50, 2.20, 5.00, 0.80, Pt(40), C_BLACK, bold=True, align=PP_ALIGN.CENTER)
    # 细分割线（装饰，非彩色块）
    add_decoration_line(s, 3.50, 3.20, 6.50, 3.20, C_DASH, 0.75)
    # 副标题
    _txt(s, "光伏 + 氢气物联网云平台 / PEM 电解水制氢数学模型", 2.00, 3.40, 6.00, 0.40, S_HEADING, C_BLACK, align=PP_ALIGN.CENTER)
    # 底部信息
    _txt(s, "周报周期：2026.07.27 — 2026.07.31 · 第二周", 2.50, 4.00, 5.00, 0.35, S_TAG, C_BLACK, align=PP_ALIGN.CENTER)
    _txt(s, "编制日期：2026.07.31", 2.50, 4.40, 5.00, 0.30, S_TAG, C_BLACK, align=PP_ALIGN.CENTER)
    add_page_number(s, n)
    return s

def build_dense_split(prs, n):
    """高密度内容页：演示溢出防护 —— 7条要点自动分页为2页"""
    items = [
        "① 物模型定义：7类设备含属性·事件·服务三要素",
        "② 设备注册与一机一密：DeviceManager封装，9台设备注册",
        "③ 数据初始化脚本：MongoDB建 thing_models/device_registry",
        "④ 后端API扩展：新增6类接口（物模型/注册表/影子读写）",
        "⑤ 桥接心跳跟踪：收到遥测更新心跳，15s巡检超时标记offline",
        "⑥ MQTT认证准备：自定义amqtt插件，默认匿名生产一键启用",
        "⑦ 前端设备管理页重做：状态卡+筛选+表格+详情抽屉",
        "验证：start_all后9台模拟设备全部online，API与影子读写通过",
    ]
    # 标题
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tx = TITLE_X_POOL[(n-1) % len(TITLE_X_POOL)]
    _txt(s, "二、设备接入框架（主线一）", tx, SAFE_TOP, 7.00, 0.40, S_HEADING, C_BLACK, bold=True)
    _txt(s, "目标：在既有链路基础上补全设备管理这一层", SAFE_LEFT, 0.74, SAFE_WIDTH, 0.38, S_HEADING, C_BLACK, bold=True)
    # 逐条渲染 + 溢出检测
    cy = 1.22
    page_items = []
    for it in items:
        ph = estimate_text_height(it, SAFE_WIDTH, 16)
        if cy + ph > SAFE_BOTTOM - 0.3:  # 预留标签/页码空间
            # 溢出 → 当前页收尾，剩余进下一页
            break
        tb = slide_textbox(s, it, SAFE_LEFT, cy, SAFE_WIDTH)
        page_items.append(it)
        cy += ph + 0.1
    add_page_number(s, n)
    # 若还有剩余 → 第二页（分页）
    remaining = items[len(page_items):]
    if remaining:
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        _txt(s2, "二、设备接入框架（续）", TITLE_X_POOL[n % len(TITLE_X_POOL)], SAFE_TOP, 7.00, 0.40, S_HEADING, C_BLACK, bold=True)
        cy2 = 1.22
        for it in remaining:
            ph = estimate_text_height(it, SAFE_WIDTH, 16)
            slide_textbox(s2, it, SAFE_LEFT, cy2, SAFE_WIDTH)
            cy2 += ph + 0.1
        add_page_number(s2, n+1)
    return s

def slide_textbox(slide, text, x, y, w):
    ph = estimate_text_height(text, w, 16)
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(ph+0.05))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.space_after = Pt(4)
    if '：' in text[:20]:
        head, rest = text.split('：', 1)
        r1 = p.add_run(); r1.text = head + '：'; r1.font.size = S_HEADING; r1.font.bold = True; r1.font.color.rgb = C_BLACK
        r2 = p.add_run(); r2.text = rest; r2.font.size = S_HEADING; r2.font.color.rgb = C_BLACK
    else:
        r = p.add_run(); r.text = text; r.font.size = S_HEADING; r.font.color.rgb = C_BLACK
    return tb

def build_summary_enhanced(prs, n):
    """结论页：原v6.1只有4行很空 → v7.0用内容扩展+状态卡增强"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _txt(s, "八、结论", TITLE_X_POOL[(n-1)%len(TITLE_X_POOL)], SAFE_TOP, 5.00, 0.40, S_HEADING, C_BLACK, bold=True)
    # 3张状态卡（替代纯文字，填充右侧空间）
    add_status_card(s, 6.20, 1.10, 3.30, 1.80, "设备接入", "9台模拟设备全部online，具备注册/认证/影子能力")
    add_status_card(s, 6.20, 3.10, 3.30, 1.80, "ML学习", "阶段一收口，三阶段排期至8/16")
    add_status_card(s, 6.20, 5.10, 3.30, 1.80, "蒸馏V2", "平均相似度0.7452，闭环全链路可行")
    # 左侧扩展文字（加粗标题+说明）
    left_paras = [
        "设备接入框架+物模型：完成落地并验证，为设备到货做准备",
        "机器学习基础系统学习：按三阶段排期，本周五讲完成收口",
        "大模型蒸馏闭环第二轮：V2真实运行，定位创意类短板",
        "下一步重点：保证ML连续性，推进MQTT实测与监控页",
    ]
    add_body_separate(s, left_paras, SAFE_LEFT, 1.10, 5.80)
    add_page_number(s, n)
    return s

def build_empty_risk(prs, n):
    """偏差与风险页：原v6.1只有表格很空 → 加装饰线+图标卡增强"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _txt(s, "六、偏差与风险", TITLE_X_POOL[(n-1)%len(TITLE_X_POOL)], SAFE_TOP, 5.00, 0.40, S_HEADING, C_BLACK, bold=True)
    # 表格区（虚线框）
    add_dashed_box(s, 0.09, 1.05, 9.65, 3.32)
    _txt(s, "（表格内容：进度偏差 / 技术风险 / 应对措施）", SAFE_LEFT, 2.50, SAFE_WIDTH, 0.60, S_HEADING, C_BLACK)
    # 下方：风险等级状态卡（填充原空白区）
    add_status_card(s, SAFE_LEFT, 4.60, 3.00, 1.50, "低风险", "主线一/ML按计划推进")
    add_status_card(s, 3.40, 4.60, 3.00, 1.50, "中风险", "蒸馏V2未正循环，需优化")
    add_status_card(s, 6.60, 4.60, 3.00, 1.50, "待跟进", "MQTT认证实测未启动")
    add_page_number(s, n)
    return s

def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W); prs.slide_height = Inches(SLIDE_H)
    # 1. 封面（空页增强：焦点留白+装饰线）
    build_cover(prs, 1)
    # 2. 高密度页（溢出防护：自动分页）
    build_dense_split(prs, 2)
    # 3. 结论页（空页增强：状态卡+扩展）
    build_summary_enhanced(prs, 4)
    # 4. 风险页（空页增强：装饰+卡片）
    build_empty_risk(prs, 5)
    # 每页溢出检测
    for i, slide in enumerate(prs.slides, 1):
        v = check_overflow(slide)
        print(f"第{i}页溢出检测: {'❌ '+str(v) if v else '✅ 无溢出'}")
    prs.save(r"D:\Desktop\大模型\v70_demo.pptx")
    print(f"已生成 v70_demo.pptx，共 {len(prs.slides.__iter__.__self__._sldIdLst)} 页")

if __name__ == "__main__":
    main()
