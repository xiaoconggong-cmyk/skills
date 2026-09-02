from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2],16), int(hex_str[2:4],16), int(hex_str[4:6],16))

def add_textbox(slide, text, x, y, w, h, style, profile, align=PP_ALIGN.LEFT):
    levels = profile['typography']['levels']
    level = levels.get(style, levels['body_main'])
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(level['size_pt'])
    run.font.bold = level.get('bold', False)
    color = profile['colors'].get('text_primary', '#000000')
    run.font.color.rgb = hex_to_rgb(color)
    return txBox

def add_rect(slide, x, y, w, h, fill_color=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape

# Academic profile
ap = {
    "canvas": {"width_inches": 10.0, "height_inches": 7.5},
    "colors": {"page_bg": "#FFFFFF", "text_primary": "#000000", "text_secondary": "#374151",
               "emphasis": "#FF0000", "label_tag": "#7030A0", "accent_primary": None},
    "typography": {"heading_family": None, "body_family": None,
        "levels": {"cover_title": {"size_pt": 40, "bold": True},
                   "section_title": {"size_pt": 24, "bold": True},
                   "subtitle": {"size_pt": 20, "bold": True},
                   "body_main": {"size_pt": 16, "bold": False},
                   "body_small": {"size_pt": 14, "bold": False},
                   "table_cell": {"size_pt": 12, "bold": False},
                   "tag_label": {"size_pt": 14, "bold": True}}},
    "decoration": {"level": 0, "decorative_lines": False},
    "density": {"level": "high", "max_body_lines": 8}
}

# Business profile
bp = {
    "canvas": {"width_inches": 13.33, "height_inches": 7.5},
    "colors": {"page_bg": "#FFFFFF", "text_primary": "#1E293B", "text_secondary": "#64748B",
               "emphasis": "#FF0000", "accent_primary": "#2563EB", "accent_secondary": "#DBEAFE"},
    "typography": {"heading_family": None, "body_family": None,
        "levels": {"cover_title": {"size_pt": 44, "bold": True},
                   "section_title": {"size_pt": 28, "bold": True},
                   "subtitle": {"size_pt": 20, "bold": True},
                   "body_main": {"size_pt": 16, "bold": False},
                   "body_small": {"size_pt": 14, "bold": False},
                   "table_cell": {"size_pt": 12, "bold": False},
                   "tag_label": {"size_pt": 14, "bold": True}}},
    "decoration": {"level": 3, "section_bars": True, "decorative_lines": True},
    "density": {"level": "mid", "max_body_lines": 6}
}

base = r"C:\Users\25124\.workbuddy\skills\ppt-multiagent\assets"

# ===== Academic Demo (4:3, 5 slides) =====
prs = Presentation()
prs.slide_width = Inches(10.0)
prs.slide_height = Inches(7.5)
w, h = 10.0, 7.5

# Slide 1: Cover
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, 0, 0, w, h, "#FFFFFF")
add_textbox(s, "深度学习异常检测研究进展", 0.5, 0.20*h, w-1, 0.22*h, 'cover_title', ap, PP_ALIGN.CENTER)
add_textbox(s, "主动视觉语言推理方法", 0.5, 0.55*h, w-1, 0.50, 'subtitle', ap, PP_ALIGN.CENTER)
add_textbox(s, "2026年3月 大组会汇报", 0.5, 0.85*h, w-1, 0.40, 'body_main', ap, PP_ALIGN.CENTER)

# Slide 2: Content with red emphasis + purple tags
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, 0, 0, w, h, "#FFFFFF")
add_textbox(s, "1. 研究背景", 0.10*w, 0.02*h, 0.80*w, 0.55, 'section_title', ap)
add_textbox(s, "异常检测是工业视觉中的核心挑战", 0, 0.10*h, w, 0.40, 'subtitle', ap)

body = "工业场景中，异常检测面临三大核心挑战：样本稀缺导致训练数据高度不平衡；异常形态多样使得单一模型难以泛化；实时性要求对推理速度提出严格约束。传统方法依赖人工特征工程，在复杂场景下准确率不足60%。"
txBox = s.shapes.add_textbox(Inches(w*0.02), Inches(0.17*h), Inches(w*0.96), Inches(0.45*h))
tf = txBox.text_frame; tf.word_wrap = True
para = tf.paragraphs[0]
emph = [
    ("样本稀缺", "#FF0000"), ("异常形态多样", "#FF0000"),
    ("实时性要求", "#FF0000"), ("准确率不足60%", "#FF0000")
]
remaining = body
for kw, clr in emph:
    idx = remaining.find(kw)
    if idx >= 0:
        r = para.add_run()
        r.text = remaining[:idx]
        r.font.size = Pt(16); r.font.color.rgb = hex_to_rgb('#000000')
        r = para.add_run()
        r.text = kw
        r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = hex_to_rgb(clr)
        remaining = remaining[idx+len(kw):]
r = para.add_run()
r.text = remaining; r.font.size = Pt(16); r.font.color.rgb = hex_to_rgb('#000000')

tags = ["Anomaly Detection", "Few-shot Learning", "Real-time Inference"]
for i, tag in enumerate(tags):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55*w+i*1.4), Inches(0.78*h), Inches(1.25), Inches(0.50))
    sh.fill.background(); sh.line.fill.background()
    pt = sh.text_frame.paragraphs[0]; pt.alignment = PP_ALIGN.CENTER
    r = pt.add_run(); r.text = tag; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = hex_to_rgb('#7030A0')

# Slide 3: 3-col comparison
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, 0, 0, w, h, "#FFFFFF")
add_textbox(s, "2. 相关方法对比", 0.05*w, 0.02*h, 0.90*w, 0.50, 'section_title', ap)
cols = [
    ("传统方法", "手工特征工程\n准确率 58.3%\n泛化能力弱"),
    ("CNN-based", "端到端特征学习\n准确率 76.1%\n需大量标注"),
    ("VLM-based (Ours)", "视觉语言联合\n准确率 89.4%\nFew-shot适应"),
]
for i, (title, body) in enumerate(cols):
    cx = 0.05*w + i * 0.30*w
    add_textbox(s, title, cx+0.10, 0.10*h, 0.28*w-0.20, 0.35, 'subtitle', ap, PP_ALIGN.CENTER)
    add_textbox(s, body, cx+0.10, 0.16*h, 0.28*w-0.20, 0.55*h, 'body_main', ap)

# Slide 4: Table
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, 0, 0, w, h, "#FFFFFF")
add_textbox(s, "3. 实验结果", 0.05*w, 0.02*h, 0.90*w, 0.50, 'section_title', ap)
tbl_sh = s.shapes.add_table(5, 5, Inches(0.05*w), Inches(0.12*h), Inches(0.50*w), Inches(0.58*h))
tbl = tbl_sh.table
for ci, hdr in enumerate(["Method", "MVTec", "VisA", "BTAD", "Avg"]):
    tbl.cell(0, ci).text = hdr
    for p in tbl.cell(0, ci).text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs: r.font.size = Pt(12); r.font.bold = True
for ri, row in enumerate([["PaDiM","85.7","72.1","78.4","78.7"],["PatchCore","87.3","74.6","80.2","80.7"],["CFA","88.9","76.3","82.1","82.4"],["Ours","94.2","88.5","90.3","91.0"]]):
    for ci, val in enumerate(row):
        tbl.cell(ri+1, ci).text = val
        for p in tbl.cell(ri+1, ci).text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs: r.font.size = Pt(12)
add_textbox(s, "关键发现", 0.58*w, 0.12*h, 0.38*w, 0.30, 'subtitle', ap)
add_textbox(s, "MVTec SOTA +5.3%\nVisA +12.2%\n推理 23 FPS\nFew-shot >87%", 0.58*w, 0.18*h, 0.38*w, 0.55*h, 'body_main', ap)

# Slide 5: Summary
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, 0, 0, w, h, "#FFFFFF")
add_textbox(s, "4. 总结与展望", 0.05*w, 0.02*h, 0.90*w, 0.50, 'section_title', ap)
txBox = s.shapes.add_textbox(Inches(0.06*w), Inches(0.12*h), Inches(0.55*w), Inches(0.75*h))
tf = txBox.text_frame; tf.word_wrap = True
for i, item in enumerate([
    "总结",
    "提出主动视觉语言推理框架，融合视觉特征与语义先验",
    "在 4 个基准数据集上取得 SOTA，平均提升 8.6%",
    "Few-shot 泛化能力显著优于纯视觉方法",
    "推理效率满足工业实时检测需求（23 FPS）",
    "",
    "后续计划",
    "扩展到 3D 点云异常检测场景",
    "探索多模态大模型的零样本异常识别",
    "与生产线系统集成，进行实地验证"
]):
    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = para.add_run(); r.text = f" {item}" if item in ("总结","后续计划") else f"  {item}" if item == "" else f"  {item}"
    bold = item in ("总结","后续计划")
    r.font.size = Pt(20 if bold else 16); r.font.bold = bold

prs.save(f"{base}/demo_academic.pptx")
print(f"Academic demo saved: {base}/demo_academic.pptx")

# ===== Business Demo (16:9, 2 slides) =====
prs2 = Presentation()
prs2.slide_width = Inches(13.33)
prs2.slide_height = Inches(7.5)
w2, h2 = 13.33, 7.5

s = prs2.slides.add_slide(prs2.slide_layouts[6])
add_rect(s, 0, 0, w2, h2, "#1E293B")
txBox = s.shapes.add_textbox(Inches(0.5), Inches(0.25*h2), Inches(w2-1), Inches(0.18*h2))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "AI-Native"; r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = hex_to_rgb('#FFFFFF')
txBox2 = s.shapes.add_textbox(Inches(0.5), Inches(0.45*h2), Inches(w2-1), Inches(0.50))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "下一代工业视觉检测平台"; r2.font.size = Pt(20); r2.font.color.rgb = hex_to_rgb('#94A3B8')
add_rect(s, 0.35*w2, 0.94*h2, 0.30*w2, 0.02, "#2563EB")
txBox3 = s.shapes.add_textbox(Inches(0.5), Inches(0.85*h2), Inches(w2-1), Inches(0.40))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run(); r3.text = "2026年7月 · 天使轮融资路演"; r3.font.size = Pt(16); r3.font.color.rgb = hex_to_rgb('#64748B')

s2 = prs2.slides.add_slide(prs2.slide_layouts[6])
add_rect(s2, 0, 0, w2, h2, "#FFFFFF")
add_rect(s2, 0, 0, w2, 0.03*h2, "#2563EB")
add_textbox(s2, "市场机遇", 0.06*w2, 0.05*h2, 0.50*w2, 0.55, 'section_title', bp)
txBox = s2.shapes.add_textbox(Inches(0.07*w2), Inches(0.14*h2), Inches(0.55*w2), Inches(0.65*h2))
tf = txBox.text_frame; tf.word_wrap = True
for i, item in enumerate(["全球工业视觉市场 2026 年达 180 亿美元，CAGR 12.7%", "AI 视觉检测渗透率仅 23%，替代空间巨大", "制造业人力成本年均上涨 8%，自动化需求迫切", "政策端：智能制造 2025 专项资金有力支撑"]):
    para = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r = para.add_run(); r.text = f"  {item}"; r.font.size = Pt(16); r.font.color.rgb = hex_to_rgb('#1E293B')

prs2.save(f"{base}/demo_business.pptx")
print(f"Business demo saved: {base}/demo_business.pptx")
print("All demos generated!")
