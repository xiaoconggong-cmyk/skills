"""
v6.0 验证脚本：生成 3 页 demo，验证内容自适应布局
- 第1页：紧凑概述型（模拟用户截图的"本周概览"）
- 第2页：标准列表型（中密度内容）
- 第3页：密集分块型（高密度内容）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ===== 常量 =====
SLIDE_W = 10.00
SLIDE_H = 7.50
SAFE_LEFT = 0.14
SAFE_RIGHT = 9.69
SAFE_TOP = 0.19
SAFE_BOTTOM = 7.30
SAFE_WIDTH = SAFE_RIGHT - SAFE_LEFT

C_BLACK  = RGBColor(0x00, 0x00, 0x00)
C_RED    = RGBColor(0xFF, 0x00, 0x00)
C_PURPLE = RGBColor(0x70, 0x30, 0xA0)
C_DASH   = RGBColor(0x9C, 0xA3, 0xAF)

S_HEADING = Pt(16)
S_TAG     = Pt(14)
S_TABLE   = Pt(12)


def add_dashed_box(slide, x, y, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.background()
    shape.line.color.rgb = C_DASH
    shape.line.width = Pt(0.5)
    shape.line.dash_style = 2
    return shape


def add_textbox_bold(slide, text, x, y, w, h, size, color, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = color
    return txBox


def add_body_textbox(slide, paragraphs, x, y, w, h):
    """paragraphs: list of list of dict"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for pi, para_runs in enumerate(paragraphs):
        p = tf.paragraphs[pi] if pi == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        for rd in para_runs:
            run = p.add_run()
            run.text = rd["text"]
            run.font.size = S_HEADING
            run.font.bold = rd.get("bold", False)
            run.font.color.rgb = rd.get("color", C_BLACK)
    return txBox


def estimate_text_height(text, width_inches, font_size_pt=16, line_height=1.15):
    chars_per_line = int(width_inches * 72 / (font_size_pt * 0.6))
    n_lines = max(1, (len(text) + chars_per_line - 1) // chars_per_line)
    line_h = font_size_pt * line_height / 72
    return n_lines * line_h + 0.2


def add_tag_label(slide, text, x, y, w=None, h=None):
    text_width = max(0.50, min(1.33, len(text) * 0.12))
    w = w or text_width
    h = h or 0.30
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = S_TAG
    run.font.bold = True
    run.font.color.rgb = C_PURPLE
    return shape


# ===================== 模式 1：紧凑概述型 =====================
def build_compact_page(slide, section_title, subtitle, paragraphs, tags=None):
    # 1. 章节标题
    title_y = SAFE_TOP
    add_textbox_bold(slide, section_title, x=1.39, y=title_y, w=5.00, h=0.40, size=S_HEADING, color=C_BLACK)

    # 2. 副标题
    subtitle_y = title_y + 0.40 + 0.15
    add_textbox_bold(slide, subtitle, x=SAFE_LEFT, y=subtitle_y, w=SAFE_WIDTH, h=0.40, size=S_HEADING, color=C_BLACK)

    # 3. 正文区域（动态计算高度）
    body_y = subtitle_y + 0.40 + 0.15
    total_text = "\n".join(["".join([r["text"] for r in p]) for p in paragraphs])
    body_h = estimate_text_height(total_text, SAFE_WIDTH, 16, 1.15)
    body_h = max(body_h, 0.50)

    add_dashed_box(slide, SAFE_LEFT - 0.05, body_y - 0.05, SAFE_WIDTH + 0.10, body_h + 0.10)
    add_body_textbox(slide, paragraphs, x=SAFE_LEFT, y=body_y, w=SAFE_WIDTH, h=body_h)

    # 4. 标签（可选，动态位置，放在内容下方）
    if tags:
        tag_y = body_y + body_h + 0.15
        if tag_y + 0.30 <= SAFE_BOTTOM:
            for i, tag_text in enumerate(tags[:3]):
                tag_x = SAFE_LEFT + i * 1.70
                add_tag_label(slide, tag_text, tag_x, tag_y)

    return slide


# ===================== 模式 2：标准列表型 =====================
def build_standard_page(slide, section_title, subtitle, paragraphs, images=None, tags=None):
    title_y = SAFE_TOP
    add_textbox_bold(slide, section_title, x=1.39, y=title_y, w=5.00, h=0.40, size=S_HEADING, color=C_BLACK)

    subtitle_y = title_y + 0.40 + 0.15
    add_textbox_bold(slide, subtitle, x=SAFE_LEFT, y=subtitle_y, w=SAFE_WIDTH, h=0.40, size=S_HEADING, color=C_BLACK)

    body_y = subtitle_y + 0.40 + 0.15
    total_text = "\n".join(["".join([r["text"] for r in p]) for p in paragraphs])
    body_h = estimate_text_height(total_text, SAFE_WIDTH, 16, 1.15)
    body_h = max(body_h, 0.50)

    add_dashed_box(slide, SAFE_LEFT - 0.05, body_y - 0.05, SAFE_WIDTH + 0.10, body_h + 0.10)
    add_body_textbox(slide, paragraphs, x=SAFE_LEFT, y=body_y, w=SAFE_WIDTH, h=body_h)

    current_y = body_y + body_h + 0.15

    if tags and current_y + 0.30 <= SAFE_BOTTOM:
        for i, tag_text in enumerate(tags[:3]):
            tag_x = SAFE_LEFT + i * 1.70
            add_tag_label(slide, tag_text, tag_x, current_y)

    return slide


# ===================== 模式 6：密集分块型 =====================
def build_dense_page(slide, section_title, blocks):
    add_textbox_bold(slide, section_title, x=1.39, y=SAFE_TOP, w=5.00, h=0.40, size=S_HEADING, color=C_BLACK)

    for block in blocks:
        bx, by, bw, bh = block["x"], block["y"], block["w"], block["h"]
        add_dashed_box(slide, bx - 0.05, by - 0.05, bw + 0.10, bh + 0.10)
        if block.get("title"):
            add_textbox_bold(slide, block["title"], x=bx, y=by, w=bw, h=0.30, size=S_HEADING, color=C_BLACK)
        if block.get("paragraphs"):
            add_body_textbox(slide, block["paragraphs"], x=bx, y=by + 0.30, w=bw, h=bh - 0.30)

    return slide


# ===================== 主入口 =====================
prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)

# ---- 第1页：紧凑概述型（模拟用户截图的"本周概览"）----
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
paragraphs1 = [
    [{"text": "生产环境部署：", "bold": True, "color": C_BLACK},
     {"text": "部署脚本 805 行，覆盖 RS 初始化、全量+增量备份、Prometheus 监控、", "bold": False, "color": C_BLACK},
     {"text": "8 维度健康检查", "bold": True, "color": C_RED},
     {"text": "，健康评分 78.6%。", "bold": False, "color": C_BLACK}],
    [{"text": "文档体系构建：", "bold": True, "color": C_BLACK},
     {"text": "设计文档、使用手册、运维手册、团队培训材料，共 ~76KB，8+7+8+8 章。", "bold": False, "color": C_BLACK}],
    [{"text": "蒸馏闭环（夜间）：", "bold": True, "color": C_BLACK},
     {"text": "teacher 生成 → student 训练 → 自动评估三阶段管线打通，评估集 v1→v2→v3，蒸馏后综合评分 +20.2%。", "bold": False, "color": C_BLACK}],
    [{"text": "专项协作：", "bold": True, "color": C_BLACK},
     {"text": "周一会议 4 事项（验收清单/传感器/物联网图层/通信对接）、周二论文大纲 6 章+初稿 ~8000 字、周三-周四 IoT 平台学习 2 天。", "bold": False, "color": C_BLACK}],
]
build_compact_page(slide1, "1. 本周概览", "三主线并行，四项核心任务 100% 完成", paragraphs1, tags=["生产部署", "文档体系"])

# ---- 第2页：标准列表型（中密度内容）----
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
paragraphs2 = [
    [{"text": "问题定义：", "bold": True, "color": C_BLACK},
     {"text": "现有方法多依赖", "bold": False, "color": C_BLACK},
     {"text": "一次性的全局感知", "bold": True, "color": C_RED},
     {"text": "，无法捕捉细粒度的语义关系，导致在复杂场景下性能显著下降。", "bold": False, "color": C_BLACK}],
    [{"text": "核心挑战：", "bold": True, "color": C_BLACK},
     {"text": "（1）样本稀缺问题：标注数据获取成本高，且分布不均；（2）", "bold": False, "color": C_BLACK},
     {"text": "多模态对齐", "bold": True, "color": C_RED},
     {"text": "困难：视觉-语言特征空间差异大；（3）推理效率瓶颈：大模型参数量大，部署成本高。", "bold": False, "color": C_BLACK}],
    [{"text": "解决方案：", "bold": True, "color": C_BLACK},
     {"text": "提出", "bold": False, "color": C_BLACK},
     {"text": "主动聚焦机制", "bold": True, "color": C_RED},
     {"text": "，通过动态注意力选择，在关键区域上集中计算资源，同时保持全局感知能力。", "bold": False, "color": C_BLACK}],
]
build_standard_page(slide2, "2. 问题与挑战", "主动视觉语言推理中的关键难题", paragraphs2, tags=["被动感知", "主动聚焦"])

# ---- 第3页：密集分块型（高密度内容）----
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
blocks = [
    {
        "title": "模块 A：数据预处理",
        "paragraphs": [
            [{"text": "输入：原始图像（224×224）+ 文本描述", "bold": False, "color": C_BLACK}],
            [{"text": "处理：图像编码器提取视觉特征，文本编码器提取语义特征，", "bold": False, "color": C_BLACK},
             {"text": "对比学习", "bold": True, "color": C_RED},
             {"text": "对齐两个空间。", "bold": False, "color": C_BLACK}],
        ],
        "x": 0.14, "y": 0.90, "w": 4.60, "h": 2.20
    },
    {
        "title": "模块 B：主动聚焦",
        "paragraphs": [
            [{"text": "机制：基于", "bold": False, "color": C_BLACK},
             {"text": "注意力热图", "bold": True, "color": C_RED},
             {"text": "动态选择 ROI，减少 60% 计算量。", "bold": False, "color": C_BLACK}],
            [{"text": "输出：聚焦区域特征 + 全局上下文向量", "bold": False, "color": C_BLACK}],
        ],
        "x": 5.20, "y": 0.90, "w": 4.60, "h": 2.20
    },
    {
        "title": "模块 C：推理融合",
        "paragraphs": [
            [{"text": "融合策略：", "bold": False, "color": C_BLACK},
             {"text": "交叉注意力", "bold": True, "color": C_RED},
             {"text": " + 门控机制，动态平衡局部与全局信息。", "bold": False, "color": C_BLACK}],
            [{"text": "性能：在 VQA 任务上提升 12.3%，推理速度提升 2.1×", "bold": False, "color": C_BLACK}],
        ],
        "x": 0.14, "y": 3.40, "w": 4.60, "h": 2.20
    },
    {
        "title": "模块 D：蒸馏优化",
        "paragraphs": [
            [{"text": "teacher-student 框架：大模型指导小模型，", "bold": False, "color": C_BLACK},
             {"text": "知识蒸馏", "bold": True, "color": C_RED},
             {"text": "保持性能同时压缩模型。", "bold": False, "color": C_BLACK}],
            [{"text": "效果：模型体积减少 75%，精度损失 < 2%", "bold": False, "color": C_BLACK}],
        ],
        "x": 5.20, "y": 3.40, "w": 4.60, "h": 2.20
    },
]
build_dense_page(slide3, "3. 系统架构", blocks)

# 保存
output = "C:\\Users\\25124\\.workbuddy\\skills\\ppt-multiagent\\assets\\demo_v6.pptx"
prs.save(output)
print(f"PPT 已生成: {output}，共 {len(prs.slides)} 页")
