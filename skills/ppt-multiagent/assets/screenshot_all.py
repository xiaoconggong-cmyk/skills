#!/usr/bin/env python3
"""
自动截图脚本：启动本地服务器，用 Playwright 打开页面并截图保存。
覆盖：
  1. demo_preview.html —— 模板效果预览页（全页截图）
  2. demo_academic.pptx —— 学术风 5 页（逐页渲染为 HTML 后截图）
  3. demo_business.pptx —— 商务风 2 页（逐页渲染为 HTML 后截图）
"""

import os
import sys
import time
import json
import http.server
import socketserver
import threading
from pathlib import Path

ASSETS_DIR = Path(r"C:\Users\25124\.workbuddy\skills\ppt-multiagent\assets")
OUTPUT_DIR = ASSETS_DIR / "screenshots"
OUTPUT_DIR.mkdir(exist_ok=True)

# ===== Step 1: Start HTTP server =====
PORT = 8765

class QuietHandler(http.server.SimpleHTTPRequestHandler):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, directory=str(ASSETS_DIR), **kwargs)
    def log_message(self, format, *args):
        pass  # suppress logs

def start_server():
    server = socketserver.TCPServer(("", PORT), QuietHandler)
    t = threading.Thread(target=server.serve_forever, daemon=True)
    t.start()
    return server

print("Starting HTTP server on port", PORT)
server = start_server()
time.sleep(0.5)

# ===== Step 2: Screenshot demo_preview.html =====
from playwright.sync_api import sync_playwright

with sync_playwright() as p:
    browser = p.chromium.launch(headless=True)
    context = browser.new_context(viewport={"width": 1400, "height": 900})

    # --- Screenshot 1: demo_preview.html (full page) ---
    print("Screenshotting demo_preview.html ...")
    page = context.new_page()
    page.goto(f"http://localhost:{PORT}/demo_preview.html", wait_until="networkidle")
    time.sleep(1)
    # Get full page height
    full_height = page.evaluate("document.body.scrollHeight")
    page.set_viewport_size({"width": 1400, "height": full_height})
    time.sleep(0.5)
    page.screenshot(path=str(OUTPUT_DIR / "01_demo_preview_full.png"), full_page=True)
    print(f"  -> {OUTPUT_DIR / '01_demo_preview_full.png'}")

    # --- Screenshot 2-3: demo_preview.html academic section ---
    print("Screenshotting demo_preview.html academic section ...")
    page.set_viewport_size({"width": 1400, "height": 900})
    page.goto(f"http://localhost:{PORT}/demo_preview.html", wait_until="networkidle")
    time.sleep(0.5)
    # Scroll to academic section
    academic_section = page.query_selector("#academic")
    if academic_section:
        academic_section.scroll_into_view_if_needed()
        time.sleep(0.5)
        page.screenshot(path=str(OUTPUT_DIR / "02_demo_academic_section.png"), full_page=False)
        print(f"  -> {OUTPUT_DIR / '02_demo_academic_section.png'}")

    # --- Screenshot 4: demo_preview.html business section ---
    print("Screenshotting demo_preview.html business section ...")
    business_section = page.query_selector("#business")
    if business_section:
        business_section.scroll_into_view_if_needed()
        time.sleep(0.5)
        page.screenshot(path=str(OUTPUT_DIR / "03_demo_business_section.png"), full_page=False)
        print(f"  -> {OUTPUT_DIR / '03_demo_business_section.png'}")

    # ===== Step 3: Render PPTX slides as HTML and screenshot =====
    from pptx import Presentation
    from pptx.util import Inches, Pt
    import base64
    from io import BytesIO

    def render_pptx_to_html(pptx_path, profile_type):
        """Convert PPTX slides to an HTML page, one slide per section."""
        prs = Presentation(pptx_path)

        slide_html_parts = []
        for idx, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            # Detect text role by position/font
                            top_pct = shape.top / prs.slide_height
                            left_pct = shape.left / prs.slide_width
                            font_size = None
                            is_bold = False
                            for run in para.runs:
                                if run.font.size:
                                    font_size = run.font.size / 12700
                                    is_bold = run.font.bold or False
                                    break
                            texts.append({
                                "text": t,
                                "top_pct": top_pct,
                                "left_pct": left_pct,
                                "font_size": font_size,
                                "is_bold": is_bold
                            })

            # Determine slide type based on position/size heuristics
            slide_type = "content"
            bg_color = "#FFFFFF"
            text_color = "#1E293B"

            if profile_type == "academic":
                if idx == 0:
                    slide_type = "cover"
                elif idx == len(prs.slides) - 1:
                    slide_type = "summary"
            elif profile_type == "business":
                if idx == 0:
                    slide_type = "cover"
                    bg_color = "#1E293B"
                    text_color = "#FFFFFF"

            # Build HTML for this slide
            html_parts = []
            ratio_w, ratio_h = (10, 7.5) if profile_type == "academic" else (13.33, 7.5)
            slide_w = 900
            slide_h = int(slide_w * ratio_h / ratio_w)

            # scale factor
            sf_w = slide_w / (ratio_w * 72)
            sf_h = slide_h / (ratio_h * 72)

            html_parts.append(f'<div class="slide" style="width:{slide_w}px;height:{slide_h}px;'
                              f'background:{bg_color};position:relative;margin:16px auto;'
                              f'box-shadow:0 2px 12px rgba(0,0,0,0.12);border-radius:4px;overflow:hidden;">')

            for tinfo in texts:
                t = tinfo["text"]
                tp = tinfo["top_pct"]
                lp = tinfo["left_pct"]
                fs = tinfo.get("font_size") or 16

                top_px = int(tp * slide_h)
                left_px = int(lp * slide_w)
                max_w = int(slide_w - left_px - 20)

                if tp < 0.05:  # section title
                    clr = text_color
                    fw = "700"
                elif tp < 0.12:  # subtitle
                    clr = text_color if profile_type == "academic" else "#94A3B8"
                    fw = "600"
                elif idx == 0 and profile_type == "business" and tp < 0.30:
                    clr = "#FFFFFF"
                    fw = "700"
                else:
                    clr = text_color if profile_type == "academic" else "#1E293B"
                    fw = "600" if tinfo.get("is_bold") else "400"

                html_parts.append(
                    f'<div style="position:absolute;top:{top_px}px;left:{left_px}px;'
                    f'max-width:{max_w}px;font-size:{fs}px;font-weight:{fw};'
                    f'color:{clr};line-height:1.5;padding:0 4px;">{t}</div>'
                )

            html_parts.append(f'<div style="position:absolute;bottom:4px;right:8px;'
                              f'font-size:10px;color:#999;">Slide {idx+1}</div>')
            html_parts.append('</div>')
            slide_html_parts.append(''.join(html_parts))

        # Wrap in full HTML doc
        doc = f'''<!DOCTYPE html>
<html lang="zh">
<head><meta charset="utf-8"><title>{profile_type.upper()} Demo</title>
<style>
body {{ font-family: "Microsoft YaHei", sans-serif; background: #f0f0f0; padding: 20px; }}
h2 {{ text-align: center; color: #333; }}
</style></head>
<body>
<h2>{profile_type.upper()} Style Demo ({len(prs.slides)} slides)</h2>
{"".join(slide_html_parts)}
</body></html>'''

        html_path = ASSETS_DIR / f"_temp_{profile_type}_slides.html"
        html_path.write_text(doc, encoding="utf-8")
        return html_path

    # --- Academic PPTX screenshots ---
    print("Rendering academic PPTX slides to HTML...")
    academic_html = render_pptx_to_html(
        str(ASSETS_DIR / "demo_academic.pptx"), "academic"
    )
    print("Screenshotting academic slides...")
    page.goto(f"http://localhost:{PORT}/{academic_html.name}", wait_until="networkidle")
    time.sleep(0.5)
    full_height = page.evaluate("document.body.scrollHeight")
    page.set_viewport_size({"width": 960, "height": full_height})
    time.sleep(0.3)
    page.screenshot(path=str(OUTPUT_DIR / "04_demo_academic_slides.png"), full_page=True)
    print(f"  -> {OUTPUT_DIR / '04_demo_academic_slides.png'}")

    # --- Business PPTX screenshots ---
    print("Rendering business PPTX slides to HTML...")
    business_html = render_pptx_to_html(
        str(ASSETS_DIR / "demo_business.pptx"), "business"
    )
    print("Screenshotting business slides...")
    page.goto(f"http://localhost:{PORT}/{business_html.name}", wait_until="networkidle")
    time.sleep(0.5)
    full_height = page.evaluate("document.body.scrollHeight")
    page.set_viewport_size({"width": 960, "height": full_height})
    time.sleep(0.3)
    page.screenshot(path=str(OUTPUT_DIR / "05_demo_business_slides.png"), full_page=True)
    print(f"  -> {OUTPUT_DIR / '05_demo_business_slides.png'}")

    # --- Individual academic slides (close-up) ---
    print("Screenshotting individual academic slides...")
    page.goto(f"http://localhost:{PORT}/{academic_html.name}", wait_until="networkidle")
    time.sleep(0.5)
    slides = page.query_selector_all(".slide")
    for i, slide_el in enumerate(slides[:3]):  # first 3
        slide_el.scroll_into_view_if_needed()
        time.sleep(0.3)
        slide_el.screenshot(path=str(OUTPUT_DIR / f"06_academic_slide_{i+1}.png"))
        print(f"  -> {OUTPUT_DIR / f'06_academic_slide_{i+1}.png'}")

    # --- Individual business slides (close-up) ---
    print("Screenshotting individual business slides...")
    page.goto(f"http://localhost:{PORT}/{business_html.name}", wait_until="networkidle")
    time.sleep(0.5)
    slides = page.query_selector_all(".slide")
    for i, slide_el in enumerate(slides[:2]):
        slide_el.scroll_into_view_if_needed()
        time.sleep(0.3)
        slide_el.screenshot(path=str(OUTPUT_DIR / f"07_business_slide_{i+1}.png"))
        print(f"  -> {OUTPUT_DIR / f'07_business_slide_{i+1}.png'}")

    browser.close()

# Cleanup temp files
for f in ASSETS_DIR.glob("_temp_*_slides.html"):
    f.unlink()

server.shutdown()
print("\n=== All screenshots saved to:", OUTPUT_DIR, "===")
print("Files:", [f.name for f in sorted(OUTPUT_DIR.iterdir())])
