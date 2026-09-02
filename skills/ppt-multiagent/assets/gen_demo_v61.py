"""
v6.1 验证脚本：生成 5 页 demo，逐项验证修复效果

改进验证项：
  1. 可变标题 X（5 页 5 个不同值）
  2. 独立 textbox（每页要点独立渲染）
  3. 虚线框按需（仅图片区/多区块加框）
  4. 页码（每页右下角）
  5. 模式 7 混排型（左文右图）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ===== 常量 =====
SLIDE_W = 10.00
SLIDE_H = 7.50
SAFE_LEFT = 0.14
SAFE_RIGHT = 9.69
SAFE_TOP = 0.19
SAFE_BOTTOM = 7.30
SAFE_WIDTH = SAFE_RIGHT - SAFE_LEFT
PAGE_NUM_X = 9.40
PAGE_NUM_Y = 7.10

C_BLACK  = RGBColor(0x00, 0x00, 0x00)
C_RED    = RGBColor(0xFF, 0x00, 0x00)
C_PURPLE = RGBColor(0x70, 0x30, 0xA0)
C_DASH   = RGBColor(0x9C, 0xA3, 0xAF)

S_HEADING  = Pt(16)
S_TAG      = Pt(14)
S_TABLE    = Pt(12)
S_PAGE_NUM = Pt(10)

TITLE_X_POOL = [1.39, 1.35, 1.28, 1.20, 1.14]

def get_title_x(slide_index):
    return TITLE_X_POOL[slide_index % len(TITLE_X_POOL)]

# ===== 辅助函数 =====

def add_textbox_bold(slide, text, x, y, w, h, size, color, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = size; run.font.bold = True; run.font.color.rgb = color
    return txBox

def estimate_text_height(text, width_inches, font_size_pt=16, line_height=1.15):
    chars_per_line = int(width_inches * 72 / (font_size_pt * 0.55))
    n_lines = max(1, (len(text) + chars_per_line - 1) // chars_per_line)
    line_h = font_size_pt * line_height / 72.0
    return n_lines * line_h + 0.15

def add_body_textboxes_separate(slide, paragraphs, x, start_y, width):
    current_y = start_y
    textboxes = []
    for pi, para_runs in enumerate(paragraphs):
        para_text = "".join([rd.get("text", "") for rd in para_runs])
        para_h = estimate_text_height(para_text, width)
        if current_y + para_h > SAFE_BOTTOM:
            break
        txBox = slide.shapes.add_textbox(
            Inches(x), Inches(current_y), Inches(width), Inches(para_h + 0.05))
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.space_after = Pt(4)
        for rd in para_runs:
            run = p.add_run()
            run.text = rd.get("text", ""); run.font.size = S_HEADING
            run.font.bold = rd.get("bold", False)
            run.font.color.rgb = rd.get("color", C_BLACK)
        textboxes.append(txBox)
        current_y += para_h + 0.05
    return current_y, textboxes

def add_dashed_box(slide, x, y, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.background(); shape.line.color.rgb = C_DASH
    shape.line.width = Pt(0.5); shape.line.dash_style = 2
    return shape

def add_page_number(slide, page_num):
    txBox = slide.shapes.add_textbox(Inches(PAGE_NUM_X), Inches(PAGE_NUM_Y), Inches(0.50), Inches(0.25))
    tf = txBox.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    run = p.add_run(); run.text = str(page_num)
    run.font.size = S_PAGE_NUM; run.font.color.rgb = C_DASH; run.font.bold = False

def add_tag_label(slide, text, x, y, w=None, h=None):
    text_width = max(0.50, min(1.33, len(text) * 0.24))  # v6.1 修正
    w = w or text_width; h = h or 0.30
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.background(); shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text; run.font.size = S_TAG; run.font.bold = True; run.font.color.rgb = C_PURPLE
    return shape

def add_image_clean(slide, img_path, x, y, w, h):
    try:
        pic = slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(w), Inches(h))
        pic.line.fill.background()
    except:
        # 图片不存在时画一个占位矩形
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        shape.line.color.rgb = C_DASH; shape.line.width = Pt(0.5)
        tf = shape.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = f"[图片 {img_path}]"; run.font.size = Pt(10); run.font.color.rgb = C_DASH
    return None


# ===== 页面生成 =====

prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)

# ==== 第 1 页：紧凑概述型（低密度，无虚线框）====
print("生成第 1 页：紧凑概述型...")
slide = prs.slides.add_slide(prs.slide_layouts[6])
sid = 0

# 标题（可变 X）
add_textbox_bold(slide, "1. 本周概览", x=get_title_x(sid), y=SAFE_TOP, w=5.00, h=0.40, size=S_HEADING, color=C_BLACK)
# 副标题
add_textbox_bold(slide, "三主线并行，四项核心任务推进", x=SAFE_LEFT, y=0.74, w=SAFE_WIDTH, h=0.40, size=S_HEADING, color=C_BLACK)
# 正文——每个要点独立 textbox（v6.1 核心改进）
paras = [
    [{"text": "1. 生产环境部署：", "bold": True, "color": C_BLACK}, {"text": "脚本 805 行，覆盖初始化+备份+监控", "bold": False, "color": C_BLACK}],
    [{"text": "2. 监控体系搭建：", "bold": True, "color": C_BLACK}, {"text": "Prometheus + Grafana，", "bold": False, "color": C_BLACK}, {"text": "告警规则 12 条", "bold": True, "color": C_RED}],
    [{"text": "3. 自动化测试：", "bold": True, "color": C_BLACK}, {"text": "覆盖率从 45% 提升至 ", "bold": False, "color": C_BLACK}, {"text": "78%", "bold": True, "color": C_RED}],
    [{"text": "4. 文档完善：", "bold": True, "color": C_BLACK}, {"text": "API 文档 32 篇，部署指南 8 篇", "bold": False, "color": C_BLACK}],
]
current_y = 1.29
current_y, _ = add_body_textboxes_separate(slide, paras, x=SAFE_LEFT, start_y=current_y, width=SAFE_WIDTH)
# 无虚线框（单正文区）
# 无标签（低密度页面省略）
add_page_number(slide, sid + 1)

# ==== 第 2 页：紧凑概述型（中密度，验证标题 X 变化，无虚线框）====
print("生成第 2 页：紧凑概述型...")
slide = prs.slides.add_slide(prs.slide_layouts[6])
sid = 1

add_textbox_bold(slide, "2. 技术方案", x=get_title_x(sid), y=SAFE_TOP, w=5.00, h=0.40, size=S_HEADING, color=C_BLACK)
add_textbox_bold(slide, "通信协议 v3 设计方案", x=SAFE_LEFT, y=0.74, w=SAFE_WIDTH, h=0.40, size=S_HEADING, color=C_BLACK)

paras = [
    [{"text": "架构选型：", "bold": True, "color": C_BLACK}, {"text": "gRPC + Protobuf，支持流式传输和双向通信", "bold": False, "color": C_BLACK}],
    [{"text": "接口定义：", "bold": True, "color": C_BLACK}, {"text": "28 个 RPC 方法，", "bold": False, "color": C_BLACK}, {"text": "平均响应 < 5ms", "bold": True, "color": C_RED}],
    [{"text": "性能优化：", "bold": True, "color": C_BLACK}, {"text": "连接池复用，QPS 从 1.2k 提升至 ", "bold": False, "color": C_BLACK}, {"text": "8.5k", "bold": True, "color": C_RED}],
    [{"text": "兼容性：", "bold": True, "color": C_BLACK}, {"text": "向下兼容 v2 协议，平滑迁移零中断", "bold": False, "color": C_BLACK}],
]
current_y = 1.29
current_y, _ = add_body_textboxes_separate(slide, paras, x=SAFE_LEFT, start_y=current_y, width=SAFE_WIDTH)
# 有标签（中密度可加）
add_tag_label(slide, "gRPC", SAFE_LEFT, current_y + 0.05)
add_tag_label(slide, "Protobuf", SAFE_LEFT + 1.70, current_y + 0.05)
add_page_number(slide, sid + 1)

# ==== 第 3 页：标准列表型 + 底部图片（验证虚线框仅包图片集群）====
print("生成第 3 页：标准列表型 + 底部图片...")
slide = prs.slides.add_slide(prs.slide_layouts[6])
sid = 2

add_textbox_bold(slide, "3. 系统架构图", x=get_title_x(sid), y=SAFE_TOP, w=5.00, h=0.40, size=S_HEADING, color=C_BLACK)
add_textbox_bold(slide, "微服务架构 v2.0 部署拓扑", x=SAFE_LEFT, y=0.74, w=SAFE_WIDTH, h=0.40, size=S_HEADING, color=C_BLACK)

paras = [
    [{"text": "API Gateway 层：Nginx + Kong，统一入口，限流熔断", "bold": False, "color": C_BLACK}],
    [{"text": "服务层：Go 微服务集群，", "bold": False, "color": C_BLACK}, {"text": "K8s 自动扩缩", "bold": True, "color": C_RED}],
    [{"text": "数据层：MySQL 主从 + Redis Cluster + ES 日志", "bold": False, "color": C_BLACK}],
]
current_y = 1.29
current_y, _ = add_body_textboxes_separate(slide, paras, x=SAFE_LEFT, start_y=current_y, width=SAFE_WIDTH)

# 底部图片集群（虚线框包围）
img_y = current_y + 0.15
img_w, img_h = 4.50, 2.00
add_image_clean(slide, "architecture.png", SAFE_LEFT, img_y, img_w, img_h)
add_image_clean(slide, "deployment.png", SAFE_LEFT + 4.70, img_y, img_w, img_h)
add_dashed_box(slide, SAFE_LEFT - 0.05, img_y - 0.05, SAFE_WIDTH + 0.10, img_h + 0.10)

add_page_number(slide, sid + 1)

# ==== 第 4 页：模式 7 混排型（左文右图，v6.1 新增）====
print("生成第 4 页：模式 7 混排型...")
slide = prs.slides.add_slide(prs.slide_layouts[6])
sid = 3

add_textbox_bold(slide, "4. 性能对比", x=get_title_x(sid), y=SAFE_TOP, w=5.00, h=0.40, size=S_HEADING, color=C_BLACK)
add_textbox_bold(slide, "v2 vs v3 协议性能对比", x=SAFE_LEFT, y=0.74, w=SAFE_WIDTH, h=0.40, size=S_HEADING, color=C_BLACK)

# 左文（55% 宽度）
paras = [
    [{"text": "QPS 提升：", "bold": True, "color": C_BLACK}, {"text": "1.2k → ", "bold": False, "color": C_BLACK}, {"text": "8.5k (7x)", "bold": True, "color": C_RED}],
    [{"text": "平均延迟：", "bold": True, "color": C_BLACK}, {"text": "23ms → ", "bold": False, "color": C_BLACK}, {"text": "4.2ms", "bold": True, "color": C_RED}],
    [{"text": "吞吐量：", "bold": True, "color": C_BLACK}, {"text": "从 80MB/s 提升至 520MB/s", "bold": False, "color": C_BLACK}],
    [{"text": "并发连接数：", "bold": True, "color": C_BLACK}, {"text": "支持 10k 长连接（增长 5x）", "bold": False, "color": C_BLACK}],
]
body_w = 5.20
current_y = 1.29
current_y, _ = add_body_textboxes_separate(slide, paras, x=SAFE_LEFT, start_y=current_y, width=body_w)

# 右图（虚线框包围）
img_x, img_y = 5.80, 1.29
img_w, img_h = 3.60, 3.00
add_image_clean(slide, "perf_chart.png", img_x, img_y, img_w, img_h)
add_dashed_box(slide, img_x - 0.05, img_y - 0.05, img_w + 0.10, img_h + 0.10)

add_page_number(slide, sid + 1)

# ==== 第 5 页：密集分块型（每区块有虚线框）====
print("生成第 5 页：密集分块型...")
slide = prs.slides.add_slide(prs.slide_layouts[6])
sid = 4

add_textbox_bold(slide, "5. 下一阶段规划", x=get_title_x(sid), y=SAFE_TOP, w=5.00, h=0.40, size=S_HEADING, color=C_BLACK)

blocks = [
    {"title": "性能优化", "x": 0.14, "y": 0.95, "w": 4.65, "h": 2.20,
     "paras": [[{"text": "• gRPC 流式压缩算法优化", "bold": False, "color": C_BLACK}],
               [{"text": "• 连接池参数 ", "bold": False, "color": C_BLACK}, {"text": "自适应调优", "bold": True, "color": C_RED}],
               [{"text": "• 零拷贝序列化方案调研", "bold": False, "color": C_BLACK}]]},
    {"title": "功能开发", "x": 5.10, "y": 0.95, "w": 4.65, "h": 2.20,
     "paras": [[{"text": "• 双向流式 RPC 接口 12 个", "bold": False, "color": C_BLACK}],
               [{"text": "• 消息队列集成 (Kafka)", "bold": False, "color": C_BLACK}],
               [{"text": "• 健康检查 + 优雅关闭", "bold": False, "color": C_BLACK}]]},
    {"title": "测试覆盖", "x": 0.14, "y": 3.45, "w": 4.65, "h": 2.20,
     "paras": [[{"text": "• 单元测试覆盖率目标 ", "bold": False, "color": C_BLACK}, {"text": "85%", "bold": True, "color": C_RED}],
               [{"text": "• 集成测试场景 50+ 个", "bold": False, "color": C_BLACK}],
               [{"text": "• 混沌工程压测计划", "bold": False, "color": C_BLACK}]]},
    {"title": "文档交付", "x": 5.10, "y": 3.45, "w": 4.65, "h": 2.20,
     "paras": [[{"text": "• 架构设计文档 v3.0", "bold": False, "color": C_BLACK}],
               [{"text": "• 接口规范 Wiki 20 篇", "bold": False, "color": C_BLACK}],
               [{"text": "• 运维手册 + FAQ", "bold": False, "color": C_BLACK}]]},
]

for block in blocks:
    bx, by, bw, bh = block["x"], block["y"], block["w"], block["h"]
    add_dashed_box(slide, bx - 0.05, by - 0.05, bw + 0.10, bh + 0.10)
    add_textbox_bold(slide, block["title"], x=bx, y=by, w=bw, h=0.30, size=S_HEADING, color=C_BLACK)
    add_body_textboxes_separate(slide, block["paras"], x=bx, start_y=by + 0.35, width=bw)

add_page_number(slide, sid + 1)

# ==== 保存 ====
output = "D:/Desktop/大模型/v61_demo.pptx"
prs.save(output)
print(f"\n✅ 验证 demo 已生成: {output}")
print("验证项：")
print("  ✓ 标题 X 变体: [1.39, 1.35, 1.28, 1.20, 1.14]（5 页 5 个不同值）")
print("  ✓ 独立 textbox: 每页要点独立渲染")
print("  ✓ 虚线框按需: 第1/2页无框，第3页仅图片区有框，第4页仅图片区有框，第5页每区块有框")
print("  ✓ 页码: 每页右下角 10pt 灰色")
print("  ✓ 模式 7: 第4页左文右图混排")
print(f"  共 {len(prs.slides)} 页")
