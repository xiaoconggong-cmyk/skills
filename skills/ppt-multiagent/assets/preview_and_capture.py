"""
用 python-pptx 读取 PPTX，生成 HTML 预览，然后用 Playwright 截图
"""
import os
import subprocess

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_PATH = r"C:\Users\25124\.workbuddy\skills\ppt-multiagent\assets\demo_v6.pptx"
OUTPUT_DIR = r"C:\Users\25124\.workbuddy\skills\ppt-multiagent\assets\screenshots_v6"

os.makedirs(OUTPUT_DIR, exist_ok=True)

prs = Presentation(PPTX_PATH)
slide_w = prs.slide_width.inches
slide_h = prs.slide_height.inches
px_per_inch = 96

for idx, slide in enumerate(prs.slides):
    html_parts = []
    html_parts.append(f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8"><style>
body {{ margin:0; padding:20px; background:#e0e0e0; display:flex; justify-content:center; align-items:center; min-height:100vh; }}
.slide {{ width:{slide_w*px_per_inch}px; height:{slide_h*px_per_inch}px; background:#FFFFFF; position:relative; box-shadow:0 4px 20px rgba(0,0,0,0.15); overflow:hidden; }}
.shape {{ position:absolute; box-sizing:border-box; }}
</style></head><body><div class="slide">""")

    for shape in slide.shapes:
        x = shape.left.inches * px_per_inch
        y = shape.top.inches * px_per_inch
        w = shape.width.inches * px_per_inch
        h = shape.height.inches * px_per_inch

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            html_parts.append(f'<div class="shape" style="left:{x}px;top:{y}px;width:{w}px;height:{h}px;background:#f0f0f0;border:1px dashed #ccc;display:flex;align-items:center;justify-content:center;font-size:12px;color:#999;">[图片]</div>')
        elif shape.has_text_frame:
            tf = shape.text_frame
            p = tf.paragraphs[0] if tf.paragraphs else None
            if p and p.runs:
                run = p.runs[0]
                font_size = run.font.size.pt if run.font.size else 16
                bold = 'font-weight:bold;' if run.font.bold else ''
                color = f'#{run.font.color.rgb}' if run.font.color and run.font.color.rgb else '#000000'
                align = 'center' if p.alignment == 1 else ('right' if p.alignment == 2 else 'left')
                text = p.text[:300]
                has_dash = False
                if hasattr(shape, 'line') and shape.line and shape.line.dash_style is not None:
                    has_dash = True
                if has_dash:
                    html_parts.append(f'<div class="shape" style="left:{x}px;top:{y}px;width:{w}px;height:{h}px;border:1px dashed #9CA3AF;"></div>')
                else:
                    html_parts.append(f'<div class="shape" style="left:{x}px;top:{y}px;width:{w}px;height:{h}px;font-size:{font_size}px;color:{color};{bold}text-align:{align};white-space:pre-wrap;overflow:hidden;font-family:Microsoft YaHei,SimSun,sans-serif;line-height:1.15;">{text}</div>')
            else:
                if hasattr(shape, 'line') and shape.line and shape.line.dash_style is not None:
                    html_parts.append(f'<div class="shape" style="left:{x}px;top:{y}px;width:{w}px;height:{h}px;border:1px dashed #9CA3AF;"></div>')
        else:
            if hasattr(shape, 'line') and shape.line and shape.line.dash_style is not None:
                html_parts.append(f'<div class="shape" style="left:{x}px;top:{y}px;width:{w}px;height:{h}px;border:1px dashed #9CA3AF;"></div>')

    html_parts.append("</div></body></html>")

    html_path = os.path.join(OUTPUT_DIR, f"slide_{idx+1}.html")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write("".join(html_parts))
    print(f"HTML saved: {html_path}")

# 用 Playwright 截图
js_code = f'''
const {{ chromium }} = require('playwright');
(async () => {{
    const browser = await chromium.launch();
    const page = await browser.newPage({{ viewport: {{ width: 1280, height: 960 }} }});
    for (let i = 1; i <= {len(prs.slides)}; i++) {{
        await page.goto('file:///{OUTPUT_DIR.replace(chr(92), "/")}/slide_' + i + '.html');
        await page.waitForTimeout(500);
        await page.screenshot({{ path: '{OUTPUT_DIR.replace(chr(92), "/")}/slide_' + i + '.png', fullPage: false }});
        console.log('Screenshot: slide_' + i + '.png');
    }}
    await browser.close();
}})();
'''
js_path = os.path.join(OUTPUT_DIR, "screenshot.js")
with open(js_path, "w", encoding="utf-8") as f:
    f.write(js_code)

print("Running Playwright screenshot...")
result = subprocess.run(["node", js_path], capture_output=True, text=True)
print(result.stdout)
if result.stderr:
    print("STDERR:", result.stderr)
print("Done!")
