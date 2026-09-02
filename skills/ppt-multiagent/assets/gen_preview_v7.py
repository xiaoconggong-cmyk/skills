# -*- coding: utf-8 -*-
"""从 v70_demo.pptx 提取形状坐标，生成 HTML 预览（用于截图验证）。"""
from pptx import Presentation
from pptx.util import Emu
import html

SLIDE_W, SLIDE_H = 10.00, 7.50
prs = Presentation(r"D:\Desktop\大模型\v70_demo.pptx")
out_dir = r"C:\Users\25124\.workbuddy\skills\ppt-multiagent\assets\previews_v7"

import os
os.makedirs(out_dir, exist_ok=True)

def color_of(run):
    try:
        if run.font.color and run.font.color.rgb:
            return "#%s" % str(run.font.color.rgb)
    except:
        pass
    return "#000000"

for idx, slide in enumerate(prs.slides, 1):
    divs = []
    for shape in slide.shapes:
        x = shape.left / 914400.0
        y = shape.top / 914400.0
        w = shape.width / 914400.0
        h = shape.height / 914400.0
        # 容器坐标(px)：按 120px/inch
        px = lambda v: v * 120
        left, top = px(x), px(y)
        width, height = px(w), px(h)
        style = f"position:absolute;left:{left}px;top:{top}px;width:{width}px;height:{height}px;"

        shape_type = str(shape.shape_type)
        # 连接线/装饰线
        if 'CONNECTOR' in shape_type:
            style += "border:none;background:linear-gradient(#9CA3AF,#9CA3AF) no-repeat center/100% 1.5px;"
            divs.append(f'<div style="{style}"></div>')
            continue
        # 自动形状（虚线框/状态卡）
        if 'AUTO_SHAPE' in shape_type or 'RECTANGLE' in shape_type:
            is_round = 'roundRect' in shape._element.xml
            if is_round:
                style += "border:1.5px solid #9CA3AF;border-radius:8px;background:#fff;"
            else:
                style += "border:1px dashed #9CA3AF;background:transparent;"
            divs.append(f'<div style="{style}"></div>')
            continue
        # 文本框
        if shape.has_text_frame and shape.text_frame.text.strip():
            txt = ""
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    col = color_of(run)
                    sz = run.font.size.pt if run.font.size else 16
                    bold = "font-weight:bold;" if run.font.bold else ""
                    txt += f'<span style="color:{col};font-size:{sz}pt;{bold}">{html.escape(run.text)}</span>'
            style += "font-family:'Microsoft YaHei',sans-serif;line-height:1.2;overflow:hidden;"
            divs.append(f'<div style="{style}">{txt}</div>')

    html_doc = f"""<!DOCTYPE html><html><head><meta charset="utf-8">
<style>body{{margin:0;background:#eee;}}
.slide{{position:relative;width:1200px;height:900px;background:#fff;margin:20px auto;box-shadow:0 2px 10px rgba(0,0,0,.2);}}
</style></head><body>
<div class="slide">{"".join(divs)}</div></body></html>"""
    with open(f"{out_dir}/slide_{idx}.html", "w", encoding="utf-8") as f:
        f.write(html_doc)
    print(f"生成 slide_{idx}.html")

print("完成")
